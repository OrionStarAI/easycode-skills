#!/usr/bin/env python3
"""Assemble slide PNGs into a 16:9 PPTX (images fill each slide).

Usage:
  python build_pptx.py --images "/tmp/sanwan_slides/slide_*.png" --out /tmp/sanwan_slides/output.pptx
  python build_pptx.py --images slide_01.png slide_02.png --out out.pptx
"""

from __future__ import annotations

import argparse
import glob
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print(
        "Error: python-pptx is required. Install with: pip install python-pptx",
        file=sys.stderr,
    )
    raise SystemExit(2)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def expand_images(patterns: list[str]) -> list[Path]:
    files: list[Path] = []
    for p in patterns:
        matches = sorted(glob.glob(p))
        if matches:
            files.extend(Path(m) for m in matches)
        else:
            path = Path(p)
            if path.is_file():
                files.append(path)
            else:
                print(f"Warning: no files matched: {p}", file=sys.stderr)
    seen = set()
    unique: list[Path] = []
    for f in files:
        rp = f.resolve()
        if rp not in seen:
            seen.add(rp)
            unique.append(rp)
    return unique


def build_pptx(images: list[Path], out_path: Path) -> Path:
    if not images:
        print("Error: no image files to assemble", file=sys.stderr)
        raise SystemExit(1)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    for img in images:
        if not img.is_file():
            print(f"Error: image not found: {img}", file=sys.stderr)
            raise SystemExit(1)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img),
            Inches(0),
            Inches(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )
        print(f"  + {img.name}", file=sys.stderr)

    out_path = out_path.expanduser().resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(str(out_path))
    return out_path


def main() -> int:
    parser = argparse.ArgumentParser(description="Build 16:9 PPTX from slide images")
    parser.add_argument(
        "--images",
        nargs="+",
        required=True,
        help='Image paths or globs, e.g. "/tmp/sanwan_slides/slide_*.png"',
    )
    parser.add_argument("--out", required=True, help="Output .pptx path")
    args = parser.parse_args()

    images = expand_images(args.images)
    images = sorted(images, key=lambda p: p.name)
    build_pptx(images, Path(args.out))
    return 0


if __name__ == "__main__":
    sys.exit(main())
